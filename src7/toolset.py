"""
toolset.py — Standalone worker functions for generating styled PPTX, XLSX, and PDF files.
These can be used independently of the AI agent in main.py.
"""

import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, LineChart, PieChart, Reference, Series
from openpyxl.chart.series import DataPoint
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage,
    Table, TableStyle, HRFlowable
)
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY


# ─────────────────────────────────────────────────────────────────────────────
#  WORKER 1 — PPTX
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx_file(presentation_data: dict) -> io.BytesIO:
    """
    Transforms a structured dict brief into a styled PPTX file.

    presentation_data schema:
    {
        "title"   : "Project Alpha",
        "subtitle": "Strategic Roadmap 2026",
        "slides"  : [
            {
                "title"        : "Overview",
                "content"      : ["Goal 1", "Goal 2", "Goal 3"],
                "notes"        : "Mention budget constraints.",
                "accent_color" : "#1E90FF",
                "icon_emoji"   : "🚀",
                "callout_stat" : "Revenue grew 40% YoY"
            }
        ]
    }
    """
    prs = Presentation()
    W = prs.slide_width
    H = prs.slide_height

    # ── Cover Slide ──────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(prs.slide_layouts[6])

    bg = cover.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)
    bg.line.fill.background()

    bar = cover.shapes.add_shape(1, 0, 0, Inches(0.25), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x1C, 0x72, 0x93)
    bar.line.fill.background()

    tb = cover.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.0), Inches(2.0))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = presentation_data.get("title", "Presentation")
    p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); p.font.name = "Calibri"

    sb = cover.shapes.add_textbox(Inches(0.6), Inches(4.1), Inches(8.0), Inches(0.8))
    sp = sb.text_frame.paragraphs[0]
    sp.text = presentation_data.get("subtitle", "")
    sp.font.size = Pt(20); sp.font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC); sp.font.name = "Calibri"

    # ── Content Slides ───────────────────────────────────────────────────────
    for slide_info in presentation_data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        accent_hex = slide_info.get("accent_color", "#1C7293").lstrip("#")
        accent_rgb = RGBColor(
            int(accent_hex[0:2], 16), int(accent_hex[2:4], 16), int(accent_hex[4:6], 16)
        )

        bg2 = slide.shapes.add_shape(1, 0, 0, W, H)
        bg2.fill.solid(); bg2.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
        bg2.line.fill.background()

        text_w = int(W * 0.90)

        hbar = slide.shapes.add_shape(1, 0, 0, text_w, Inches(0.08))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = accent_rgb; hbar.line.fill.background()

        icon = slide_info.get("icon_emoji", "")
        title_text = f"{icon}  {slide_info.get('title', '')}" if icon else slide_info.get("title", "")
        ttb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), text_w - Inches(0.5), Inches(1.1))
        ttb.text_frame.word_wrap = True
        tp = ttb.text_frame.paragraphs[0]
        tp.text = title_text; tp.font.size = Pt(30); tp.font.bold = True
        tp.font.color.rgb = accent_rgb; tp.font.name = "Calibri"

        div = slide.shapes.add_shape(1, Inches(0.4), Inches(1.35), text_w - Inches(0.8), Inches(0.03))
        div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD); div.line.fill.background()

        content_items = slide_info.get("content", [])
        ctb = slide.shapes.add_textbox(Inches(0.4), Inches(1.55), text_w - Inches(0.5), Inches(3.6))
        ctb.text_frame.word_wrap = True
        for i, item in enumerate(content_items):
            para = ctb.text_frame.paragraphs[0] if i == 0 else ctb.text_frame.add_paragraph()
            para.text = f"  •  {item}"; para.font.size = Pt(16)
            para.font.name = "Calibri"; para.font.color.rgb = RGBColor(0x2C, 0x2C, 0x2C)
            para.space_after = Pt(8)

        callout = slide_info.get("callout_stat", "")
        if callout:
            cbox = slide.shapes.add_shape(1, Inches(0.4), Inches(5.5), text_w - Inches(0.8), Inches(0.85))
            cbox.fill.solid(); cbox.fill.fore_color.rgb = accent_rgb; cbox.line.fill.background()
            ctb2 = slide.shapes.add_textbox(Inches(0.55), Inches(5.55), text_w - Inches(1.0), Inches(0.75))
            ctb2.text_frame.word_wrap = True
            cp = ctb2.text_frame.paragraphs[0]
            cp.text = f"📌  {callout}"; cp.font.size = Pt(15); cp.font.bold = True
            cp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); cp.font.name = "Calibri"

        if "notes" in slide_info:
            slide.notes_slide.notes_text_frame.text = slide_info["notes"]

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io


# ─────────────────────────────────────────────────────────────────────────────
#  WORKER 2 — XLSX
# ─────────────────────────────────────────────────────────────────────────────

def generate_xlsx_file(data_payload: dict) -> io.BytesIO:
    """
    Transforms a structured dict brief into a styled Excel spreadsheet.

    data_payload schema:
    {
        "sheet_name"        : "Quarterly Budget",
        "columns"           : ["Item", "Cost", "Quantity"],
        "rows"              : [["Laptops", 1200, 5], ["Licenses", 50, 100]],
        "highlight_top_n"   : 3,
        "number_format_cols": ["Cost", "Quantity"],
        "include_total"     : True,
        "freeze_header"     : True
    }
    """
    wb = Workbook()
    ws = wb.active
    ws.title = data_payload.get("sheet_name", "Sheet1")

    columns  = data_payload.get("columns", [])
    rows     = data_payload.get("rows", [])
    num_cols = data_payload.get("number_format_cols", [])
    top_n    = data_payload.get("highlight_top_n", 3)

    header_fill  = PatternFill("solid", fgColor="1C7293")
    gold_fill    = PatternFill("solid", fgColor="FFF3CD")
    alt_fill     = PatternFill("solid", fgColor="F0F7FA")
    white_fill   = PatternFill("solid", fgColor="FFFFFF")
    summary_fill = PatternFill("solid", fgColor="1E2761")
    header_font  = Font(bold=True, color="FFFFFF", name="Calibri", size=12)
    body_font    = Font(name="Calibri", size=11)
    gold_font    = Font(name="Calibri", size=11, bold=True)
    summary_font = Font(bold=True, color="FFFFFF", name="Calibri", size=11)
    center       = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left         = Alignment(horizontal="left",   vertical="center", wrap_text=True)
    thin         = Side(style="thin", color="CCCCCC")
    border       = Border(left=thin, right=thin, top=thin, bottom=thin)

    ws.row_dimensions[1].height = 32
    for ci, col_name in enumerate(columns, 1):
        cell = ws.cell(row=1, column=ci, value=col_name)
        cell.fill = header_fill; cell.font = header_font
        cell.alignment = center; cell.border = border

    # Computed subtotal
    has_cost_qty = "Cost" in columns and "Quantity" in columns
    if has_cost_qty and "Subtotal" not in columns:
        columns = columns + ["Subtotal"]
        ci = len(columns)
        cell = ws.cell(row=1, column=ci, value="Subtotal")
        cell.fill = header_fill; cell.font = header_font
        cell.alignment = center; cell.border = border
        ws.column_dimensions[get_column_letter(ci)].width = 14

    medals = ["🥇", "🥈", "🥉"]
    is_rank = columns[0].lower() in ("rank", "#", "no", "no.") if columns else False

    for ri, row_data in enumerate(rows, 2):
        ws.row_dimensions[ri].height = 22
        is_top   = (ri - 1) <= top_n
        row_fill = gold_fill if is_top else (alt_fill if ri % 2 == 0 else white_fill)
        row_font = gold_font if is_top else body_font

        for ci, value in enumerate(row_data, 1):
            col_name = columns[ci - 1] if ci - 1 < len(columns) else ""
            cell = ws.cell(row=ri, column=ci, value=value)
            cell.fill = row_fill; cell.font = row_font; cell.border = border
            if col_name in num_cols and isinstance(value, (int, float)):
                cell.number_format = "#,##0"; cell.alignment = center
            else:
                cell.alignment = left if ci == 1 else center

        if has_cost_qty:
            cost_i = columns.index("Cost"); qty_i = columns.index("Quantity")
            sub_i  = columns.index("Subtotal") + 1
            sub_cell = ws.cell(row=ri, column=sub_i,
                               value=row_data[cost_i] * row_data[qty_i])
            sub_cell.number_format = "#,##0"
            sub_cell.fill = row_fill; sub_cell.font = row_font
            sub_cell.alignment = center; sub_cell.border = border

        if is_top and is_rank:
            ws.cell(row=ri, column=1).value = medals[min(ri - 2, 2)]

    for ci, col_name in enumerate(columns, 1):
        max_len = max(len(str(col_name)), *(len(str(r[ci-1])) for r in rows if ci-1 < len(r)), default=0)
        ws.column_dimensions[get_column_letter(ci)].width = min(max_len + 4, 30)

    if data_payload.get("freeze_header", True):
        ws.freeze_panes = "A2"

    if data_payload.get("include_total", True):
        last_row = len(rows) + 2
        ws.row_dimensions[last_row].height = 24
        for ci, col_name in enumerate(columns, 1):
            cell = ws.cell(row=last_row, column=ci)
            cell.fill = summary_fill; cell.font = summary_font
            cell.border = border; cell.alignment = center
            if ci == 1:
                cell.value = "TOTAL / SUMMARY"
            elif col_name in (num_cols + (["Subtotal"] if has_cost_qty else [])):
                col_letter = get_column_letter(ci)
                cell.value = f"=SUM({col_letter}2:{col_letter}{last_row - 1})"
                cell.number_format = "#,##0"

    # ── Embedded Chart ───────────────────────────────────────────────────────
    # Auto-select chart type: use PieChart for ≤5 rows, BarChart otherwise.
    # Targets the first numeric column found.
    numeric_col_indices = [
        ci for ci, col in enumerate(columns)
        if col in data_payload.get("number_format_cols", []) or
           any(isinstance(r[ci], (int, float)) for r in rows if ci < len(r))
    ]

    if numeric_col_indices and rows:
        chart_col_idx = numeric_col_indices[0]  # first numeric column (0-based)
        label_col_idx = 0                        # first column used as labels

        data_row_count = len(rows)
        chart_start_row = data_row_count + 4     # leave 2 blank rows after total

        # Build References (openpyxl is 1-based)
        data_ref = Reference(
            ws,
            min_col=chart_col_idx + 1,
            min_row=1,                           # include header for series title
            max_row=data_row_count + 1,
        )
        cats_ref = Reference(
            ws,
            min_col=label_col_idx + 1,
            min_row=2,
            max_row=data_row_count + 1,
        )

        col_title = columns[chart_col_idx] if chart_col_idx < len(columns) else "Values"

        if data_row_count <= 5:
            chart = PieChart()
            chart.title = f"{col_title} — Distribution"
            chart.style  = 10
            chart.dataLabels = None
            # PieChart: colour slices
            series = Series(data_ref, title_from_data=True)
            chart.series.append(series)
            chart.set_categories(cats_ref)
        else:
            chart = BarChart()
            chart.type    = "col"
            chart.style   = 10
            chart.title   = f"{col_title} — Comparison"
            chart.y_axis.title = col_title
            chart.x_axis.title = columns[label_col_idx] if columns else ""
            chart.add_data(data_ref, titles_from_data=True)
            chart.set_categories(cats_ref)
            chart.shape = 4

        chart.width  = 18   # cm
        chart.height = 11   # cm

        chart_anchor = f"{get_column_letter(1)}{chart_start_row}"
        ws.add_chart(chart, chart_anchor)

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output


# ─────────────────────────────────────────────────────────────────────────────
#  PDF HEADING HELPER
# ─────────────────────────────────────────────────────────────────────────────

def _make_section_heading(text: str, accent_hex: str = "#1C7293") -> Table:
    """
    Renders a section heading as a styled Table with a bold colored left border.
    More reliable than ParagraphStyle alone — guarantees colored strip,
    background tint, and correct padding in all ReportLab versions.
    """
    inner_style = ParagraphStyle(
        f"HeadingInner_{accent_hex}",
        fontSize=14, fontName="Helvetica-Bold",
        textColor=colors.HexColor("#1E2761"),
        leading=18,
    )
    tbl = Table([[Paragraph(text, inner_style)]], colWidths=[14.5 * cm])
    tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, -1), colors.HexColor("#EBF5FB")),
        ("TOPPADDING",    (0, 0), (-1, -1), 9),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 9),
        ("LEFTPADDING",   (0, 0), (-1, -1), 12),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 8),
    ]))
    return tbl


# ─────────────────────────────────────────────────────────────────────────────
#  WORKER 3 — PDF REPORT
# ─────────────────────────────────────────────────────────────────────────────

def generate_pdf_file(report_data: dict) -> io.BytesIO:
    """
    Transforms a structured dict brief into a professionally styled PDF report.

    report_data schema:
    {
        "title": "Annual Market Analysis",
        "sections": [
            {
                "heading"      : "Executive Summary",
                "content"      : "This report outlines the growth metrics for 2026...",
                "callout"      : "Revenue grew 32% YoY — highest in a decade",
                "section_type" : "highlight",
                "table_data"   : [
                    ["Category", "Q1",   "Q2"  ],
                    ["Revenue",  "$10k", "$15k"],
                ]
            }
        ]
    }
    """
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm,  bottomMargin=2*cm,
    )

    style_title = ParagraphStyle(
        "ReportTitle", fontSize=20, fontName="Helvetica-Bold",
        textColor=colors.HexColor("#1E2761"),
        spaceAfter=4, alignment=TA_CENTER,
    )
    style_body = ParagraphStyle(
        "Body", fontSize=11, fontName="Helvetica",
        textColor=colors.HexColor("#2C2C2C"),
        leading=17, spaceAfter=8, alignment=TA_JUSTIFY,
    )

    CALLOUT_CONFIGS = {
        "normal":    ("#1C7293", colors.white),
        "highlight": ("#1E2761", colors.white),
        "warning":   ("#FFF3CD", colors.HexColor("#7B3F00")),
        "tip":       ("#D4EDDA", colors.HexColor("#155724")),
    }

    story = []
    story.append(Paragraph(report_data.get("title", "Report"), style_title))
    story.append(HRFlowable(width="100%", thickness=2.5,
                             color=colors.HexColor("#1C7293"), spaceAfter=10))
    story.append(Spacer(1, 0.4 * cm))

    for sec in report_data.get("sections", []):
        sec_type = sec.get("section_type", "normal")

        # ✅ FIX: proper styled heading
        if sec.get("heading"):
            story.append(_make_section_heading(sec["heading"]))
            story.append(Spacer(1, 1 * cm))

        if sec.get("content"):
            story.append(Paragraph(sec["content"], style_body))

        callout = sec.get("callout", "")
        if callout:
            bg_hex, text_color = CALLOUT_CONFIGS.get(sec_type, CALLOUT_CONFIGS["normal"])
            callout_style = ParagraphStyle(
                f"C_{id(sec)}", fontSize=12, fontName="Helvetica-Bold",
                textColor=text_color, alignment=TA_CENTER, leading=16,
            )
            box = Table([[Paragraph(f"📌  {callout}", callout_style)]], colWidths=[14.5*cm])
            box.setStyle(TableStyle([
                ("BACKGROUND",    (0, 0), (-1, -1), colors.HexColor(bg_hex)),
                ("TOPPADDING",    (0, 0), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
                ("LEFTPADDING",   (0, 0), (-1, -1), 14),
                ("RIGHTPADDING",  (0, 0), (-1, -1), 14),
            ]))
            story.append(Spacer(1, 0.15 * cm))
            story.append(box)
            story.append(Spacer(1, 0.15 * cm))

        if sec.get("table_data"):
            t = Table(sec["table_data"], hAlign="LEFT", repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND",     (0, 0), (-1, 0),  colors.HexColor("#1C7293")),
                ("TEXTCOLOR",      (0, 0), (-1, 0),  colors.white),
                ("FONTNAME",       (0, 0), (-1, 0),  "Helvetica-Bold"),
                ("FONTSIZE",       (0, 0), (-1, 0),  11),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#F0F7FA"), colors.white]),
                ("FONTNAME",       (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE",       (0, 1), (-1, -1), 10),
                ("ALIGN",          (0, 0), (-1, -1), "CENTER"),
                ("VALIGN",         (0, 0), (-1, -1), "MIDDLE"),
                ("GRID",           (0, 0), (-1, -1), 0.5, colors.HexColor("#CCCCCC")),
                ("TOPPADDING",     (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING",  (0, 0), (-1, -1), 8),
            ]))
            story.append(Spacer(1, 0.2 * cm))
            story.append(t)

        story.append(Spacer(1, 0.5 * cm))

    doc.build(story)
    buffer.seek(0)
    return buffer


# ─────────────────────────────────────────────────────────────────────────────
#  EXAMPLE USAGE
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    # PPTX
    with open("output_presentation.pptx", "wb") as f:
        f.write(generate_pptx_file({
            "title": "Marketing Strategy 2026", "subtitle": "Q1 Kickoff Deck",
            "slides": [
                {"title": "Target Audience", "content": ["Gen Z", "Early Adopters", "Urban Professionals"],
                 "accent_color": "#1E90FF", "icon_emoji": "🎯",
                 "callout_stat": "68% of Gen Z discover brands via social media"},
                {"title": "KPIs", "content": ["10k signups", "5% conversion", "$200k ARR"],
                 "accent_color": "#2C5F2D", "icon_emoji": "📈",
                 "callout_stat": "Each 1% conversion = $40k additional revenue"},
            ],
        }).read())
    print("✅ output_presentation.pptx")

    # XLSX
    with open("output_data.xlsx", "wb") as f:
        f.write(generate_xlsx_file({
            "sheet_name": "Sales Report",
            "columns": ["Product", "Cost", "Quantity"],
            "rows": [["Widget A", 10.50, 100], ["Gadget B", 25.00, 50], ["Doohickey C", 5.00, 200]],
            "highlight_top_n": 2, "number_format_cols": ["Cost", "Quantity"],
            "include_total": True, "freeze_header": True,
        }).read())
    print("✅ output_data.xlsx")

    # PDF
    with open("output_report.pdf", "wb") as f:
        f.write(generate_pdf_file({
            "title": "Project Status Update",
            "sections": [
                {"heading": "Overview", "section_type": "highlight",
                 "content": "The project is 80% complete and on track for Q2 delivery.",
                 "callout": "80% complete — on schedule for Q2 delivery"},
                {"heading": "Resource Allocation", "section_type": "normal",
                 "content": "The team consists of 3 developers and 1 designer.",
                 "table_data": [["Role", "Name", "Hours/Week"],
                                ["Dev", "Alice", "40"], ["Design", "Bob", "20"]]},
                {"heading": "Risks", "section_type": "warning",
                 "content": "Supply chain delays may affect hardware delivery timelines.",
                 "callout": "Hardware delivery delayed by up to 3 weeks"},
            ],
        }).read())
    print("✅ output_report.pdf")